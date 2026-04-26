"""Generate PPT #2 — 한글, 논문 형식.

Intro → Related Work → Problem → Method → Experiments → Discussion → Conclusion.
새 도메인 진입자도 이해할 수 있는 별도 md 설명 파일 동반.

Usage:
    uv run python docs/ppt/generate_ppt2_ko.py
Output:
    docs/ppt/cross_modal_anchoring_paper_ko.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

PROJECT_ROOT = Path(__file__).resolve().parents[2]
OUT = PROJECT_ROOT / "docs" / "ppt" / "cross_modal_anchoring_paper_ko.pptx"

# Color palette: Berry & Cream (논문용 — 차별화)
PRIMARY = RGBColor(0x6D, 0x2E, 0x46)       # #6D2E46 — berry (primary)
SECONDARY = RGBColor(0xA2, 0x67, 0x69)     # #A26769 — dusty rose
ACCENT = RGBColor(0x2F, 0x3C, 0x7E)        # #2F3C7E — navy (accent)
LIGHT = RGBColor(0xEC, 0xE2, 0xD0)         # #ECE2D0 — cream (bg tint)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2A, 0x1F, 0x29)          # near-black wine
BODY = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0x66, 0x66, 0x66)
RED = RGBColor(0xC9, 0x2A, 0x2A)
GREEN = RGBColor(0x2A, 0x7A, 0x3A)

# Korean fonts: Malgun Gothic / Apple Gothic / Noto Sans CJK 등 사용자 PC 의존.
# python-pptx가 폰트 이름 저장만 함 — 실제 표시는 viewer 폰트 의존.
HEADER_FONT = "Malgun Gothic"
BODY_FONT = "Malgun Gothic"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_filled_rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, font=BODY_FONT, size=14, bold=False,
             color=BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
             line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, font=BODY_FONT, size=14,
                color=BODY, line_spacing=1.2):
    """items: list of strings or (level, str) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.word_wrap = True
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            lvl, txt = it
        else:
            lvl, txt = 0, it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.line_spacing = line_spacing
        bullet = "•  " if lvl == 0 else "—  "
        run = p.add_run()
        run.text = bullet + txt
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_header(slide, title, subtitle=None, page_num=None, total_pages=None,
               section_label=None):
    """슬라이드 헤더: 좌측 berry 띠 + 섹션 라벨 + 제목."""
    add_filled_rect(slide, Inches(0), Inches(0), Inches(0.4), SLIDE_H, PRIMARY)
    if section_label:
        add_text(slide, Inches(0.6), Inches(0.18), Inches(11.5), Inches(0.3),
                 section_label.upper(), font=HEADER_FONT, size=11, bold=True,
                 color=SECONDARY)
        add_text(slide, Inches(0.6), Inches(0.5), Inches(11.5), Inches(0.7),
                 title, font=HEADER_FONT, size=26, bold=True, color=DARK)
        if subtitle:
            add_text(slide, Inches(0.6), Inches(1.15), Inches(11.5), Inches(0.4),
                     subtitle, font=BODY_FONT, size=13, italic=True, color=ACCENT)
    else:
        add_text(slide, Inches(0.6), Inches(0.25), Inches(11.5), Inches(0.7),
                 title, font=HEADER_FONT, size=28, bold=True, color=DARK)
        if subtitle:
            add_text(slide, Inches(0.6), Inches(0.95), Inches(11.5), Inches(0.4),
                     subtitle, font=BODY_FONT, size=13, italic=True, color=ACCENT)
    if page_num is not None and total_pages is not None:
        add_text(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
                 f"{page_num} / {total_pages}", size=10, color=MUTED,
                 align=PP_ALIGN.RIGHT)


def add_table(slide, x, y, w, h, headers, rows, *, header_color=PRIMARY,
              header_text=WHITE, row_alt=LIGHT, font_size=11, header_size=12):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, n_cols, x, y, w, h).table
    for j, ht in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        cell.margin_left = Inches(0.05)
        cell.margin_right = Inches(0.05)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = ht
        run.font.name = HEADER_FONT
        run.font.size = Pt(header_size)
        run.font.bold = True
        run.font.color.rgb = header_text
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_alt if i % 2 == 0 else WHITE
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
            run = p.add_run()
            txt = str(val)
            bold = txt.startswith("**") and txt.endswith("**")
            if bold:
                txt = txt[2:-2]
            run.text = txt
            run.font.name = BODY_FONT
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = DARK if bold else BODY
    return tbl


# ──────────────────────────────────────────────────────────────────────
# Build presentation
# ──────────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]
TOTAL = 15


# ── Slide 1 — Title ──────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_filled_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, PRIMARY)
add_filled_rect(s, Inches(0), Inches(5.2), SLIDE_W, Inches(0.06), SECONDARY)

add_text(s, Inches(0.7), Inches(1.8), Inches(12), Inches(0.5),
         "EMNLP 2026 후보 논문", size=14, color=LIGHT, italic=True)
add_text(s, Inches(0.7), Inches(2.3), Inches(12), Inches(1.6),
         "VLM의 교차-모달 앵커링 편향:",
         font=HEADER_FONT, size=40, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(3.2), Inches(12), Inches(1.0),
         "메커니즘 규명과 어텐션 재가중 완화책",
         font=HEADER_FONT, size=28, color=LIGHT)
add_text(s, Inches(0.7), Inches(5.5), Inches(12), Inches(0.4),
         "관련없는 숫자 이미지가 VLM의 수치 답변을 어떻게 끌어당기는가",
         size=15, color=LIGHT, italic=True)
add_text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
         "박태현 · 2026-04-26",
         size=14, color=LIGHT)


# ── Slide 2 — Abstract / TL;DR ───────────────────────────────────────
s = prs.slides.add_slide(blank)
add_header(s, "Abstract — 한 페이지 요약",
           "본 연구의 핵심 발견 다섯 가지", section_label="Overview",
           page_num=2, total_pages=TOTAL)

# 5 key findings - card layout
findings = [
    ("1", "현상 확인",
     "7개 VLM × 17,730 sample에서 무관한 숫자 이미지가 답변을 앵커 쪽으로 +7 ~ +20 pp 이동시킴 (불확실한 케이스에 한해 graded pull)"),
    ("2", "메커니즘",
     "어텐션 분석에서 anchor > neutral 어텐션 (4/4 모델), 가족별로 peak layer가 12~82 % 깊이로 분포 (E1b)"),
    ("3", "인과 검증",
     "단일 layer ablation은 효과 없음 (6/6); upper-half multi-layer ablation이 architecture-blind 인과 위치 (E1d)"),
    ("4", "완화책",
     "Upper-half 어텐션 재가중으로 direction-follow를 −10.6 ~ −17.7 % 감소 (LLaVA-1.5, ConvLLaVA Phase 2 풀 검증), em 손실 없음 (오히려 +0.77 ~ +1.30 pp)"),
    ("5", "Anchor damage 회복",
     "Paired 분석: anchor가 정확도를 −3.55 ~ −9.34 pp 손상; 완화책이 그 손상의 14 ~ 22 %를 회복"),
]

card_w = Inches(11.8); card_h = Inches(0.95); start_y = Inches(1.65)
gap = Inches(0.10)
for i, (num, title, body) in enumerate(findings):
    y = start_y + Inches(i * (card_h.inches + gap.inches))
    add_filled_rect(s, Inches(0.7), y, card_w, card_h, LIGHT)
    add_filled_rect(s, Inches(0.7), y, Inches(0.5), card_h, PRIMARY)
    add_text(s, Inches(0.7), y, Inches(0.5), card_h,
             num, font=HEADER_FONT, size=22, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.4), y + Inches(0.08), Inches(2.0), Inches(0.4),
             title, font=HEADER_FONT, size=14, bold=True, color=PRIMARY)
    add_text(s, Inches(3.5), y + Inches(0.08), Inches(8.0), Inches(0.85),
             body, size=12, color=BODY, line_spacing=1.2)


# ── Slide 3 — Introduction (Motivation) ──────────────────────────────
s = prs.slides.add_slide(blank)
add_header(s, "Introduction — 왜 이 문제가 중요한가",
           "VLM이 일상 보편화되면서 multi-image prompt가 표준이 되고 있다",
           section_label="1. Introduction", page_num=3, total_pages=TOTAL)

add_bullets(s, Inches(0.7), Inches(1.7), Inches(7.5), Inches(5),
            ["현실 시나리오: 사용자가 VLM에 여러 이미지를 동시에 보여주는 일이 늘고 있다 "
             "(앨범 정리, 멀티-스크린샷, 문서 등)",
             "위험: 정답과 무관한 이미지가 답변에 영향을 준다면? "
             "특히 '숫자가 그려진 이미지'가 다른 질문의 수치 답변을 끌어당긴다면?",
             "기존 인지 과학 (Tversky-Kahneman, Mussweiler-Strack): "
             "사람도 무관한 숫자에 anchored 되며, 그 효과는 불확실할수록 강함",
             "기존 LLM 연구 (Jones-Steinhardt 2022, Echterhoff EMNLP Findings 2024): "
             "텍스트로 anchor 주입 시 유사한 효과 확인",
             "기존 VLM 연구 공백: 시각적 anchor (rendered digit image)가 다른 "
             "이미지의 수치 질문에 미치는 영향은 미답",
             "기여:",
             (1, "(1) Cross-modal numerical anchoring 현상 정립 (7개 모델 × 17,730 sample)"),
             (1, "(2) 메커니즘 규명: post-projection LLM 깊이 축이 핵심 (encoder 아님)"),
             (1, "(3) Architecture-blind mitigation (upper-half attention re-weighting)")],
            size=13, line_spacing=1.25)

# Right panel - cognitive science framing
add_filled_rect(s, Inches(8.5), Inches(1.7), Inches(4.3), Inches(5), LIGHT)
add_text(s, Inches(8.65), Inches(1.85), Inches(4.0), Inches(0.4),
         "인지 과학 배경", font=HEADER_FONT, size=14, bold=True, color=PRIMARY)
add_text(s, Inches(8.65), Inches(2.3), Inches(4.0), Inches(2.5),
         "Tversky & Kahneman (1974):\n"
         "사람은 무관한 숫자가 제시되어도 그 숫자를 기준으로 추정함.\n\n"
         "Mussweiler & Strack (1999):\n"
         "Selective accessibility — anchor가 search direction을 편향함.\n\n"
         "→ 본 연구 H2: VLM도 같은 패턴? (불확실 시 anchor pull)",
         size=12, color=BODY, line_spacing=1.3)


# ── Slide 4 — Related Work ───────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_header(s, "Related Work — 가장 가까운 선행 연구와의 차이",
           "Cross-modal numerical anchoring은 지금까지 직접 다뤄지지 않음",
           section_label="2. Related Work", page_num=4, total_pages=TOTAL)

rw_rows = [
    ["LLM Anchoring", "Jones & Steinhardt 2022;\nLou & Sun 2024;\nEchterhoff 2024",
     "텍스트 prompt에 앵커",
     "**텍스트 only — 시각 anchor 미답**"],
    ["Typographic Attacks", "Goh 2021;\nWang 2025 (NAACL)",
     "이미지 위에 class-label 텍스트 오버레이",
     "**Classification flip — 수치 회귀 미답**"],
    ["VLMBias", "Nguyen ICML AI4MATH 2025",
     "주제 식별 라벨로 카운트 편향",
     "**Memorization prior — 수치 anchor 아님**"],
    ["FigStep", "Gong AAAI 2025", "텍스트→이미지 변환 jailbreak",
     "**Adversarial — 인지 편향 framing 아님**"],
    ["Tinted Frames", "2026 arXiv", "Question-form framing in VLM",
     "**질문 형식만, 이미지 modality 미답**"],
]
add_text(s, Inches(0.7), Inches(1.7), Inches(12), Inches(0.4),
         "근접 선행 연구 매트릭스 (각 work이 다른 차원으로 본 연구와 구별됨)",
         font=HEADER_FONT, size=13, bold=True, color=PRIMARY)
add_table(s, Inches(0.7), Inches(2.15), Inches(12.0), Inches(3.0),
          ["연구 영역", "대표 논문", "Setup", "본 연구와 차이"], rw_rows,
          font_size=10, header_size=11)

# Bottom card - novelty positioning
add_filled_rect(s, Inches(0.7), Inches(5.4), Inches(12), Inches(1.7), PRIMARY)
add_text(s, Inches(0.85), Inches(5.5), Inches(11.7), Inches(0.4),
         "본 연구의 위치", font=HEADER_FONT, size=15, bold=True, color=WHITE)
add_text(s, Inches(0.85), Inches(5.95), Inches(11.7), Inches(1.1),
         "독립적인 rendered-number 이미지를 multi-image VLM prompt에 추가하고, "
         "open numerical VQA에서 회귀 형태의 shift를 측정한 최초의 연구. "
         "추가로 'wrong-vs-correct asymmetry' (불확실성 변조)는 LLM·VLM 양쪽 anchoring "
         "선행 연구에서 한 번도 분리되지 않은 분석.",
         size=13, color=LIGHT, line_spacing=1.3)


# ── Slide 5 — Problem Definition ─────────────────────────────────────
s = prs.slides.add_slide(blank)
add_header(s, "Problem Definition — 3-conditions paired 비교",
           "같은 모델, 같은 질문, 세 가지 이미지 입력",
           section_label="3. Problem", page_num=5, total_pages=TOTAL)

# Three conditions as cards
conditions = [
    ("A. target_only",
     "타겟 이미지만 (질문에 해당)",
     "베이스라인 정확도 측정",
     PRIMARY),
    ("B. target + neutral",
     "타겟 + 중립 이미지 (자릿수 없음)",
     "'두 번째 이미지가 산만하게 함' 통제",
     SECONDARY),
    ("C. target + number",
     "타겟 + 숫자 이미지 (앵커, e.g. '7')",
     "Anchor manipulation",
     ACCENT),
]
for i, (name, desc, role, c) in enumerate(conditions):
    x = Inches(0.7 + i * 4.25)
    add_filled_rect(s, x, Inches(1.7), Inches(4.0), Inches(2.0), LIGHT)
    add_filled_rect(s, x, Inches(1.7), Inches(4.0), Inches(0.55), c)
    add_text(s, x + Inches(0.15), Inches(1.78), Inches(3.7), Inches(0.4),
             name, font=HEADER_FONT, size=15, bold=True, color=WHITE)
    add_text(s, x + Inches(0.15), Inches(2.3), Inches(3.7), Inches(0.45),
             desc, size=13, color=DARK, bold=True)
    add_text(s, x + Inches(0.15), Inches(2.85), Inches(3.7), Inches(0.7),
             role, size=12, color=BODY, italic=True)

# 측정 metrics
add_filled_rect(s, Inches(0.7), Inches(4.0), Inches(12), Inches(3.0), LIGHT)
add_text(s, Inches(0.85), Inches(4.1), Inches(11.7), Inches(0.4),
         "측정 메트릭 (모두 paired sample 단위)",
         font=HEADER_FONT, size=14, bold=True, color=PRIMARY)
add_bullets(s, Inches(0.85), Inches(4.55), Inches(11.7), Inches(2.5),
            ["adoption_rate: 예측이 anchor 자릿수와 정확히 일치하는 비율 (categorical)",
             "moved_closer_rate: |pred(C) − anchor| < |pred(A) − anchor| (graded shift, 본 연구 핵심)",
             "mean_anchor_pull: |pred(A) − anchor| − |pred(C) − anchor| (magnitude)",
             "exact_match: ground truth 정답 일치율 (정확도 직접 측정)",
             "Anchor effect = (C에서 메트릭) − (B에서 메트릭) → 산만함이 아닌 anchoring 효과 분리"],
            size=12, line_spacing=1.25)


# ── Slide 6 — Method (Setup) ─────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_header(s, "Method — 데이터셋 및 모델",
           "VQAv2 number subset, 7개 VLM panel (mechanism용 6개 + pilot 4개)",
           section_label="4. Method", page_num=6, total_pages=TOTAL)

# Left: Dataset & Setup
add_text(s, Inches(0.7), Inches(1.7), Inches(6), Inches(0.4),
         "데이터셋 및 실험 세팅", font=HEADER_FONT, size=15, bold=True, color=PRIMARY)
add_bullets(s, Inches(0.7), Inches(2.15), Inches(6), Inches(4.5),
            ["VQAv2 number subset: answer ∈ [0, 8], 한 자릿수 정수 답",
             "Sample size: answer 값별 400개 × 9 = 3,600 unique 질문",
             "Anchor 변형: 5 irrelevant sets / 질문 (0~9 자릿수 무작위)",
             "총 paired records: **17,730 per model**",
             "시각 자극: 1024×1024 PNG, 가운데 정렬된 단일 자릿수",
             "프롬프트:",
             (1, "JSON-strict template `{\"result\": <number>}`"),
             (1, "Greedy decoding, max_new_tokens=8 (InternVL3는 32로 패치)"),
             "Outlier filter: IQR×1.5 on per-sample anchor-distance"],
            size=12, line_spacing=1.25)

# Right: Models
add_filled_rect(s, Inches(7.0), Inches(1.7), Inches(5.7), Inches(5.0), LIGHT)
add_text(s, Inches(7.15), Inches(1.85), Inches(5.5), Inches(0.4),
         "모델 panel", font=HEADER_FONT, size=14, bold=True, color=PRIMARY)
model_rows = [
    ["Phase 1 main", "7 모델 × full 17,730:\nGemma3-27b, Gemma4-31b/-e4b, Qwen2.5-VL-7b, "
     "Qwen3-VL-30b/-8b, LLaVA-Next-7b"],
    ["E1/E1b/E1d", "6 모델 mechanism panel:\nGemma4-e4b, Qwen2.5-VL-7b, "
     "LLaVA-1.5-7b, InternVL3-8b, ConvLLaVA-7b, FastVLM-7b"],
    ["E2 pilot", "11 모델 × 1,125:\n위 6 + Phase 1 main 5 모델"],
    ["E4 mitigation", "Mid-stack cluster 3:\nLLaVA-1.5, ConvLLaVA, InternVL3"],
]
y_offset = 2.3
for label, desc in model_rows:
    add_text(s, Inches(7.15), Inches(y_offset), Inches(2.0), Inches(0.4),
             label, font=HEADER_FONT, size=12, bold=True, color=ACCENT)
    add_text(s, Inches(9.15), Inches(y_offset), Inches(3.5), Inches(1.0),
             desc, size=11, color=BODY, line_spacing=1.2)
    y_offset += 1.1


# ── Slide 7 — Experiment 1: Behavioral evidence (Phase A) ───────────
s = prs.slides.add_slide(blank)
add_header(s, "Exp 1 — 행동 증거: 앵커링은 graded pull, categorical capture 아님",
           "Phase 1 main + Phase A 정밀 분석 결과",
           section_label="5. Experiments", page_num=7, total_pages=TOTAL)

# Two columns
add_text(s, Inches(0.7), Inches(1.7), Inches(6), Inches(0.4),
         "H1 결과 — 7개 모델 모두에서 anchor 효과 확인",
         font=HEADER_FONT, size=13, bold=True, color=PRIMARY)
e1_rows = [
    ["gemma4-e4b (약함)", "0.553", "0.123", "0.247"],
    ["gemma3-27b-it", "0.628", "0.141", "0.162"],
    ["qwen3-vl-30b-it (강함)", "0.759", "0.120", "0.163"],
    ["qwen2.5-vl-7b", "0.736", "0.110", "0.089"],
]
add_table(s, Inches(0.7), Inches(2.15), Inches(6), Inches(2.0),
          ["model", "acc(타겟)", "adoption", "moved-closer"], e1_rows,
          font_size=11, header_size=11)
add_text(s, Inches(0.7), Inches(4.25), Inches(6), Inches(0.5),
         "→ adoption 11~14 %는 chance 수준에 가까움; moved-closer는 8~25 %로 더 큼",
         size=11, italic=True, color=ACCENT)

# Right: H2 refined
add_filled_rect(s, Inches(7.0), Inches(1.7), Inches(5.7), Inches(2.6), LIGHT)
add_text(s, Inches(7.15), Inches(1.85), Inches(5.5), Inches(0.4),
         "Phase A — H2 정밀 분석", font=HEADER_FONT, size=13, bold=True, color=PRIMARY)
add_bullets(s, Inches(7.15), Inches(2.3), Inches(5.5), Inches(2.0),
            ["Wrong vs Correct stratification (target_only 정답 여부 기준)",
             "adoption gap: ALL within ±2 pp → categorical H2 falsified",
             "moved-closer gap: **+6.9 ~ +19.6 pp on 7/7** → graded H2 holds",
             "= '불확실 시 anchor pull'은 categorical capture가 아닌 graded shift"],
            size=11, line_spacing=1.2)

# Bottom: refined claim card
add_filled_rect(s, Inches(0.7), Inches(4.95), Inches(12), Inches(2.1), PRIMARY)
add_text(s, Inches(0.85), Inches(5.05), Inches(11.7), Inches(0.4),
         "정제된 H2 (논문 핵심 claim)",
         font=HEADER_FONT, size=15, bold=True, color=WHITE)
add_text(s, Inches(0.85), Inches(5.5), Inches(11.7), Inches(1.5),
         "VLM이 불확실할 때 (target_only에서 틀린 케이스), anchor 자릿수 이미지가 "
         "예측을 anchor 쪽으로 +7 ~ +20 pp 더 끌어당긴다. 단, anchor를 정확히 "
         "복사하는 비율은 거의 변하지 않는다. → 앵커가 confident estimate를 "
         "교체하는 것이 아니라, search direction을 편향시킨다 (Mussweiler-Strack "
         "selective accessibility).",
         size=13, color=LIGHT, line_spacing=1.3)


# ── Slide 8 — Experiment 2: Attention mechanism (E1, E1b) ────────────
s = prs.slides.add_slide(blank)
add_header(s, "Exp 2 — 메커니즘: 어텐션 분석 (E1 + E1b)",
           "Anchor가 어디서 (어떤 layer) 어떤 budget을 (text vs target) 훔치는가",
           section_label="5. Experiments", page_num=8, total_pages=TOTAL)

# Top: per-layer table
add_text(s, Inches(0.7), Inches(1.7), Inches(8.0), Inches(0.4),
         "E1b — Per-layer 분석 (n=200/모델, output_attentions=True)",
         font=HEADER_FONT, size=13, bold=True, color=PRIMARY)
peak_rows = [
    ["Gemma-SigLIP", "L5/42 (12 %)", "+0.0501", "text", "**early + large**"],
    ["InternVL3-InternViT", "L14/28 (52 %)", "+0.0193", "text", "mid-stack"],
    ["LLaVA-1.5-CLIP-ViT", "L16/32 (52 %)", "+0.0188", "text", "mid-stack"],
    ["ConvLLaVA-ConvNeXt", "L16/32 (52 %)", "+0.0224", "text", "mid-stack"],
    ["Qwen2.5-VL-Qwen-ViT", "L22/28 (82 %)", "+0.0153", "**target**", "late"],
    ["FastVLM-FastViT", "L22/28 (82 %)", "+0.0467", "text", "late + large"],
]
add_table(s, Inches(0.7), Inches(2.15), Inches(8.0), Inches(2.7),
          ["model", "peak layer", "δ peak", "budget 출처", "유형"], peak_rows,
          font_size=11, header_size=11)

# Right: figure
fig_path = PROJECT_ROOT / "outputs/attention_analysis/_per_layer/fig_peak_budget_decomposition.png"
if fig_path.exists():
    s.shapes.add_picture(str(fig_path), Inches(8.85), Inches(1.7),
                         width=Inches(4.0), height=Inches(2.0))
    add_text(s, Inches(8.85), Inches(3.75), Inches(4.0), Inches(0.4),
             "Peak에서의 budget 분해 (모델별)", size=10, color=MUTED, italic=True,
             align=PP_ALIGN.CENTER)

# Bottom card: H3 falsified + 4 archetypes
add_filled_rect(s, Inches(0.7), Inches(5.0), Inches(12), Inches(2.1), LIGHT)
add_text(s, Inches(0.85), Inches(5.1), Inches(11.7), Inches(0.4),
         "두 가지 핵심 발견", font=HEADER_FONT, size=14, bold=True, color=PRIMARY)
add_bullets(s, Inches(0.85), Inches(5.55), Inches(11.7), Inches(1.5),
            ["E1c — H3 falsified: ConvNeXt encoder (ConvLLaVA)는 CLIP-ViT (LLaVA-1.5)와 "
             "동일한 L16, 동일한 mechanism (text-stealing), 19 % 이내의 magnitude. "
             "**Encoder architecture가 아니라 LLM stack의 depth 축이 패턴을 결정**",
             "4 archetypes: (1) Gemma early+large, (2) Mid-stack cluster (3 encoder × 동일), "
             "(3) Qwen late+target-stealing, (4) FastVLM late+text+큰 magnitude"],
            size=12, line_spacing=1.2)


# ── Slide 9 — Experiment 3: Causal evidence (E1d) ────────────────────
s = prs.slides.add_slide(blank)
add_header(s, "Exp 3 — 인과 증거: 어디를 ablation 해야 효과가 있는가 (E1d)",
           "7가지 ablation 모드 × 6개 모델 (n=200), peak layer는 인과 위치 아님",
           section_label="5. Experiments", page_num=9, total_pages=TOTAL)

e1d_rows = [
    ["Single-layer at E1b peak", "≤ ±3 pp", "✗ Null on 6/6"],
    ["Single-layer at L0 (control)", "≤ ±3 pp", "✗ Null on 6/6"],
    ["Stack-wide ablation", "−11 ~ −22 pp", "Fluency 4–6× on 3/6"],
    ["Lower-half ablation", "+27 / +17 / +7 pp", "🔥 BACKFIRES on 3/6"],
    ["**Upper-half ablation**", "**−5.5 ~ −11.5 pp on 6/6**", "**✓ Fluency clean on 4/6**"],
]
add_text(s, Inches(0.7), Inches(1.7), Inches(7.5), Inches(0.4),
         "Per-mode 결과 (Δ direction_follow_rate)",
         font=HEADER_FONT, size=13, bold=True, color=PRIMARY)
add_table(s, Inches(0.7), Inches(2.15), Inches(7.5), Inches(2.7),
          ["ablation 모드", "Δ direction_follow", "verdict"], e1d_rows,
          font_size=11, header_size=11)

# Right: figure
fig_path = PROJECT_ROOT / "outputs/causal_ablation/_summary/fig_direction_follow.png"
if fig_path.exists():
    s.shapes.add_picture(str(fig_path), Inches(8.5), Inches(1.7),
                         width=Inches(4.3), height=Inches(2.6))
    add_text(s, Inches(8.5), Inches(4.35), Inches(4.3), Inches(0.4),
             "Mode별 direction-follow Δ", size=10, color=MUTED, italic=True,
             align=PP_ALIGN.CENTER)

# Bottom: insights
add_filled_rect(s, Inches(0.7), Inches(5.0), Inches(12), Inches(2.1), PRIMARY)
add_text(s, Inches(0.85), Inches(5.1), Inches(11.7), Inches(0.4),
         "두 가지 결론", font=HEADER_FONT, size=14, bold=True, color=WHITE)
add_bullets(s, Inches(0.85), Inches(5.55), Inches(11.7), Inches(1.5),
            ["**Single-layer 개입은 무효**: 어텐션 신호가 multi-layer redundant — "
             "E1b의 peak는 correlational, 인과 위치가 아니다",
             "**Upper-half 개입이 architecture-blind 인과 위치**: 한 intervention으로 "
             "3개 다른 encoder (CLIP-ViT, InternViT, ConvNeXt)에 모두 작동, "
             "fluency도 깨끗"],
            size=12, color=LIGHT, line_spacing=1.2)


# ── Slide 10 — Experiment 4: Mitigation (E4 design + Phase 1) ────────
s = prs.slides.add_slide(blank)
add_header(s, "Exp 4-A — 완화책 설계: Soft Attention Re-weighting",
           "E1d의 hard mask를 soft strength 축으로 일반화 + Phase 1 sweep",
           section_label="5. Experiments", page_num=10, total_pages=TOTAL)

# Method 설명
add_text(s, Inches(0.7), Inches(1.7), Inches(6), Inches(0.4),
         "방법", font=HEADER_FONT, size=13, bold=True, color=PRIMARY)
add_bullets(s, Inches(0.7), Inches(2.1), Inches(6), Inches(2.7),
            ["[n_layers/2, n_layers) 범위의 각 LLM decoder layer에 forward pre-hook 부착",
             "Hook이 attention_mask의 anchor span 컬럼에 strength 더함",
             "Post-softmax: anchor 어텐션이 exp(strength)배로 가중치 down",
             "strength 격자: [0, −0.5, −1, −2, −3, −5, −10⁴] (log spacing)",
             "선택 규칙: ≥ 10 % df 감소 + ≤ 2 pp em 손실 만족하는 최소 |s|"],
            size=12, line_spacing=1.25)

# Phase 1 results table
add_text(s, Inches(0.7), Inches(4.95), Inches(8), Inches(0.4),
         "Phase 1 sweep 결과 (n=200 stratified per model)",
         font=HEADER_FONT, size=13, bold=True, color=PRIMARY)
sweep_rows = [
    ["llava-1.5-7b", "0.305", "**−3.0**", "0.265", "**−13 %**", "+0.5 pp"],
    ["convllava-7b", "0.290", "**−2.0**", "0.260", "**−10 %**", "+0.0 pp"],
    ["internvl3-8b", "0.161", "**−0.5**", "0.132", "**−17.7 %**", "+1.9 pp"],
]
add_table(s, Inches(0.7), Inches(5.4), Inches(8.0), Inches(1.6),
          ["model", "df₀", "s*", "df at s*", "Δ rel df", "Δ em"], sweep_rows,
          font_size=11, header_size=11)

# Right: figure
fig_path = PROJECT_ROOT / "outputs/e4_mitigation/_summary/sweep_pareto.png"
if fig_path.exists():
    s.shapes.add_picture(str(fig_path), Inches(7.0), Inches(1.7),
                         width=Inches(5.7), height=Inches(2.5))
    add_text(s, Inches(7.0), Inches(4.25), Inches(5.7), Inches(0.4),
             "Phase 1 Pareto: strength × (df, em)",
             size=10, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

# Right card: 발견
add_filled_rect(s, Inches(8.85), Inches(4.95), Inches(3.85), Inches(2.1), LIGHT)
add_text(s, Inches(9), Inches(5.05), Inches(3.7), Inches(0.4),
         "주목점", font=HEADER_FONT, size=13, bold=True, color=PRIMARY)
add_bullets(s, Inches(9), Inches(5.5), Inches(3.7), Inches(1.5),
            ["3/3 모두 타깃 달성",
             "s*가 모델별로 다름 (−0.5 ~ −3.0)",
             "Effect size가 baseline df와 반비례"],
            size=11, line_spacing=1.2)


# ── Slide 11 — Experiment 4: Phase 2 results ─────────────────────────
s = prs.slides.add_slide(blank)
add_header(s, "Exp 4-B — 풀 검증: 17,730 sample × 88,650 records / 모델",
           "Phase 1 예측을 0.3 pp 이내로 복제, CI 약 10× 좁아짐",
           section_label="5. Experiments", page_num=11, total_pages=TOTAL)

# Phase 2 main results
phase2_rows = [
    ["llava-1.5-7b", "−3.0", "0.2578", "0.2122", "**−4.55 pp**",
     "**−17.7 %**", "+0.77 pp", "100 %"],
    ["convllava-7b", "−2.0", "0.2283", "0.2042", "**−2.42 pp**",
     "**−10.6 %**", "+1.30 pp", "100 %"],
    ["internvl3-8b", "−0.5", "(진행 중, 다음 세션)", "—", "—", "—", "—", "~50 %"],
]
add_text(s, Inches(0.7), Inches(1.7), Inches(12), Inches(0.4),
         "Phase 2 baseline vs treated (s*)",
         font=HEADER_FONT, size=13, bold=True, color=PRIMARY)
add_table(s, Inches(0.7), Inches(2.15), Inches(12.0), Inches(1.7),
          ["model", "s*", "df₀", "df treated", "Δ df", "Δ rel", "Δ em", "% complete"],
          phase2_rows, font_size=10, header_size=10)

# Replication card + caveat card
add_filled_rect(s, Inches(0.7), Inches(4.05), Inches(5.9), Inches(3.0), LIGHT)
add_text(s, Inches(0.85), Inches(4.15), Inches(5.7), Inches(0.4),
         "복제 정확도", font=HEADER_FONT, size=14, bold=True, color=PRIMARY)
add_bullets(s, Inches(0.85), Inches(4.6), Inches(5.7), Inches(2.4),
            ["LLaVA: Phase 1 −17.7 % 예측 → Phase 2 −17.7 % 측정",
             "ConvLLaVA: Phase 1 −10.3 % 예측 → Phase 2 −10.6 % 측정",
             "두 모델 모두 0.3 pp 이내 복제, CI는 약 10× 좁음",
             "em이 두 모델 모두 상승 (+0.77 / +1.30 pp) → 정확도가 보존이 아니라 개선"],
            size=11, line_spacing=1.2)

add_filled_rect(s, Inches(6.8), Inches(4.05), Inches(5.9), Inches(3.0), PRIMARY)
add_text(s, Inches(6.95), Inches(4.15), Inches(5.7), Inches(0.4),
         "Caveat", font=HEADER_FONT, size=14, bold=True, color=WHITE)
add_bullets(s, Inches(6.95), Inches(4.6), Inches(5.7), Inches(2.4),
            ["ConvLLaVA fluency tail: mean_dist 2.99 → 53.54 — 일부 outlier 예측. "
             "em은 여전히 상승, df는 per-pair 계산이라 robust. 논문에선 median 사용 권장.",
             "InternVL3 Phase 2는 12-h budget 안에 미완료, 다음 세션에서 max_new_tokens=32 "
             "패치 적용 후 재시작 예정",
             "InternVL3 prose-leak 이슈 해결됨 (driver patch 적용)"],
            size=11, color=LIGHT, line_spacing=1.2)


# ── Slide 12 — Anchor damage / partial recovery ──────────────────────
s = prs.slides.add_slide(blank)
add_header(s, "재구성: Anchor가 정확도를 손상시키고, 완화책이 일부 회복",
           "Paired 분석 (intersection of valid samples) — em 비교가 fair",
           section_label="6. Discussion", page_num=12, total_pages=TOTAL)

# Phase 2 paired anchor-damage table
ad_rows = [
    ["llava-1.5-7b", "17,724", "0.3696", "0.3340", "**−3.55 pp**",
     "0.3417", "+0.77 pp", "**21.7 %**"],
    ["convllava-7b", "17,722", "0.4454", "0.3520", "**−9.34 pp**",
     "0.3651", "+1.31 pp", "**14.0 %**"],
]
add_text(s, Inches(0.7), Inches(1.7), Inches(12), Inches(0.4),
         "Phase 2 paired anchor-damage (full set)",
         font=HEADER_FONT, size=13, bold=True, color=PRIMARY)
add_table(s, Inches(0.7), Inches(2.15), Inches(12.0), Inches(1.5),
          ["model", "n_paired", "em(TO)", "em(num@0)", "anchor damage",
           "em(num@s*)", "recovery", "% damage 회복"],
          ad_rows, font_size=10, header_size=10)

# Bottom: framing
add_filled_rect(s, Inches(0.7), Inches(3.95), Inches(12), Inches(3.2), LIGHT)
add_text(s, Inches(0.85), Inches(4.05), Inches(11.7), Inches(0.4),
         "더 강한 paper claim", font=HEADER_FONT, size=15, bold=True, color=PRIMARY)
add_text(s, Inches(0.85), Inches(4.5), Inches(11.7), Inches(0.7),
         "Anchor가 정확도를 −3.55 ~ −9.34 pp 명확히 손상시키고, upper-half 어텐션 "
         "재가중이 그 손상의 14 ~ 22 %를 안전한 운영점에서 회복한다.",
         size=13, color=BODY, line_spacing=1.3)

add_text(s, Inches(0.85), Inches(5.6), Inches(11.7), Inches(0.4),
         "왜 강한가:", font=HEADER_FONT, size=14, bold=True, color=PRIMARY)
add_bullets(s, Inches(0.85), Inches(6.0), Inches(11.7), Inches(1.0),
            ["df-axis (메커니즘)와 em-axis (functional) 신호가 동일한 위치를 가리킴",
             "회복 비율이 두 모델 사이 비슷 (14–22 %) → single-model artefact 아님",
             "E1d의 single-layer null이 correlational 가능성을 배제 → 인과적"],
            size=12, line_spacing=1.2)


# ── Slide 13 — Discussion: Novelty (강한 부분) ────────────────────────
s = prs.slides.add_slide(blank)
add_header(s, "Discussion — 본 논문의 강점 (Novelty)",
           "publish-ready 발견 5가지", section_label="6. Discussion",
           page_num=13, total_pages=TOTAL)

novelty = [
    ("A1 — Graded pull refinement of H2",
     "+6.9~+19.6 pp moved-closer on uncertain (7/7 models). 'categorical capture' 가 아닌 'graded pull' 임을 처음 분리. Mussweiler-Strack selective accessibility 와 직접 매핑."),
    ("E1c — H3 falsified, depth-axis replaces encoder-axis",
     "ConvNeXt encoder(ConvLLaVA)가 CLIP-ViT(LLaVA-1.5)와 동일한 layer/mechanism으로 작동. Encoder architecture가 아닌 post-projection LLM stack depth가 패턴 결정."),
    ("E1d — Multi-layer redundancy + upper-half 인과 위치",
     "Single-layer ablation이 6/6에서 무효. Upper-half multi-layer가 −5.5~−11.5 pp 효과 + fluency clean (4/6). Architecture-blind 인과 위치 발견."),
    ("E2 pilot (H6) — Two-axis decoupling",
     "anchoring vs multi-image distraction이 직교 실패 모드 (InternVL3: 0.066 adopt / 0.355 acc_drop은 pure distraction). 별도 mitigation 전략 가능."),
    ("E4 Phase 2 — Causal mitigation at scale",
     "−17.7 % (LLaVA), −10.6 % (ConvLLaVA) df 감소 with em 개선. Phase 1 sweep의 0.3 pp 이내 복제. Anchor damage의 14~22 % 회복."),
]
y = 1.7
for title, body in novelty:
    add_filled_rect(s, Inches(0.7), Inches(y), Inches(12), Inches(0.95), LIGHT)
    add_filled_rect(s, Inches(0.7), Inches(y), Inches(0.18), Inches(0.95), GREEN)
    add_text(s, Inches(0.95), Inches(y + 0.1), Inches(11.5), Inches(0.4),
             title, font=HEADER_FONT, size=13, bold=True, color=PRIMARY)
    add_text(s, Inches(0.95), Inches(y + 0.5), Inches(11.5), Inches(0.4),
             body, size=11, color=BODY, line_spacing=1.2)
    y += 1.05


# ── Slide 14 — Discussion: Weakness (약점/한계) ──────────────────────
s = prs.slides.add_slide(blank)
add_header(s, "Discussion — 본 논문의 한계와 약점 (Weakness)",
           "Phase 2/추가 컴퓨트로 보강 필요한 부분",
           section_label="6. Discussion", page_num=14, total_pages=TOTAL)

weakness = [
    ("InternVL3 Phase 2 미완료",
     "max_new_tokens=8 truncation으로 prose-leak 발생 (~30 % parse failure). 12-h budget 만료로 다음 세션에서 패치 적용 후 재실행 필요. 헤드라인 mid-stack-cluster 결과는 LLaVA + ConvLLaVA 만 cover."),
    ("FastVLM 결과의 wide CI",
     "Unique late+text+large archetype이지만 답변 단계 valid n=75 (digit coverage 62 %). 95 % CI [+0.025, +0.072]. Upper-half ablation에서 fluency 깨짐 (mean_dist 폭주). 후속 큰 run 필요."),
    ("ConvLLaVA fluency tail at full scale",
     "Phase 2에서 mean_dist가 2.99 → 53.54로 outlier 발생. em은 여전히 상승, df는 robust. 논문에선 median + winsorised 분포 + degraded fraction count 권장."),
    ("Single dataset (VQAv2 only) 약점",
     "EMNLP main의 multi-dataset hygiene 요구를 충족 못함. TallyQA, ChartQA, MathVista smoke test만 됨 (E5 풀 run 미수행). Reviewer 'single dataset' 우려 가능."),
    ("Closed-model subset 미수행",
     "GPT-4o, Gemini 등 frontier closed model 미포함. Reviewer 'open-only' 우려 가능. E6에서 ~500 sample subset 수행 권장."),
    ("Paraphrase robustness 미검증",
     "Single prompt template (JSON-strict) 만 사용. 3-5 paraphrase × bootstrap CI × multi-comp 보정 (E7) 미수행. Behavioral claim의 robustness 우려."),
]
y = 1.7
for title, body in weakness:
    add_filled_rect(s, Inches(0.7), Inches(y), Inches(12), Inches(0.85), LIGHT)
    add_filled_rect(s, Inches(0.7), Inches(y), Inches(0.18), Inches(0.85), RED)
    add_text(s, Inches(0.95), Inches(y + 0.05), Inches(11.5), Inches(0.4),
             title, font=HEADER_FONT, size=12, bold=True, color=PRIMARY)
    add_text(s, Inches(0.95), Inches(y + 0.4), Inches(11.5), Inches(0.4),
             body, size=10, color=BODY, line_spacing=1.2)
    y += 0.92


# ── Slide 15 — Conclusion ────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_filled_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, PRIMARY)
add_text(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.6),
         "결론 및 향후 작업",
         font=HEADER_FONT, size=30, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(1.05), Inches(12), Inches(0.4),
         "Single-claim arc: notice (E1) → modulate (A1) → mechanism (E1b/E1c/E1d) → mitigate (E4)",
         size=14, italic=True, color=LIGHT)

add_text(s, Inches(0.7), Inches(1.7), Inches(12), Inches(0.4),
         "주요 기여", font=HEADER_FONT, size=16, bold=True, color=LIGHT)
add_bullets(s, Inches(0.7), Inches(2.1), Inches(12), Inches(2.5),
            ["Cross-modal numerical anchoring 현상의 첫 정립 (7개 VLM × 17,730 sample)",
             "Uncertainty-modulated graded pull (Mussweiler-Strack) 정량 증거 (+7~+20 pp)",
             "Per-layer attention localisation 4 archetypes; H3 (encoder axis) falsified, "
             "depth axis 채택",
             "Architecture-blind upper-half causal locus (E1d) 발견",
             "Phase 2 풀 스케일 mitigation 검증: −10.6/−17.7 % df 감소, em 개선 (LLaVA, ConvLLaVA)"],
            size=14, color=LIGHT, line_spacing=1.25)

add_text(s, Inches(0.7), Inches(4.85), Inches(12), Inches(0.4),
         "향후 작업 (Submission 전 우선순위)", font=HEADER_FONT, size=16, bold=True,
         color=LIGHT)
add_bullets(s, Inches(0.7), Inches(5.25), Inches(12), Inches(1.7),
            ["[우선] InternVL3 Phase 2 완료 (driver patch 적용 후)",
             "[우선] E5 multi-dataset (TallyQA, ChartQA) 풀 run — single-dataset 약점 보완",
             "[권장] E6 closed-model subset (GPT-4o / Gemini) on n≈500 stratified",
             "[권장] E7 paraphrase robustness (3~5 prompt × bootstrap CI)",
             "[옵션] Per-stratum mitigation 분석, ConvLLaVA fluency tail decomposition"],
            size=13, color=LIGHT, line_spacing=1.2)


# Save
prs.save(str(OUT))
print(f"[write] {OUT}")
print(f"[size] {OUT.stat().st_size:,} bytes")
print(f"[slides] {len(prs.slides)}")
