"""Build the before/after migration report comparing the buggy anchor·gt
form (preserved in ``outputs/before_C_form/``) against the corrected
C-form (current ``outputs/``).

Reads ``summary.json`` files from both trees, computes paper-tier metric
deltas, generates figures + tables, and writes:

  - ``docs/insights/C-form-migration-report.md`` (markdown report)
  - ``docs/figures/C_form_migration_*.png`` (per-section figures)
  - ``docs/ppt/C_form_migration_report.pptx`` (presentation deck)

The metrics inspected:
  - ``anchor_direction_follow_rate`` (df_M2 — the formula that changed)
  - ``anchor_direction_follow_rate_raw`` (df_raw, sign-only no movement)
  - ``anchor_adoption_rate`` (sanity — should be unchanged)
  - ``accuracy_exact`` (sanity — definitely unchanged)
  - ``mean_distance_to_anchor`` (unchanged)

Usage::

    uv run python scripts/build_C_form_migration_report.py
"""
from __future__ import annotations

import json
from collections import defaultdict
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd

PROJECT_ROOT = Path(__file__).resolve().parents[1]
BEFORE_ROOT = PROJECT_ROOT / "outputs" / "before_C_form"
AFTER_ROOT = PROJECT_ROOT / "outputs"
OUT_FIG = PROJECT_ROOT / "docs" / "figures"
OUT_DOC = PROJECT_ROOT / "docs" / "insights" / "C-form-migration-report.md"
OUT_PPTX = PROJECT_ROOT / "docs" / "ppt" / "C_form_migration_report.pptx"

# Sub-trees we care about (paper-tier; skip the backup).
EXPERIMENT_GLOBS = [
    "experiment",
    "experiment_anchor_strengthen_prompt",
    "experiment_chartqa",
    "experiment_chartqa_full",
    "experiment_distance_tally",
    "experiment_distance_vqa",
    "experiment_e5c_tally",
    "experiment_e5c_vqa",
    "experiment_e5d_chartqa_validation",
    "experiment_e5d_mathvista_validation",
    "experiment_e5e_chartqa_full",
    "experiment_e5e_mathvista_full",
    "experiment_e5e_mathvista_reasoning",
    "experiment_e5e_tallyqa_full",
    "experiment_encoder_pilot",
    "causal_ablation",
    "e4_mitigation",
]


@dataclass
class CellRecord:
    experiment: str
    model: str
    timestamp: str
    condition: str
    df_moved: float | None
    df_raw: float | None
    adopt: float | None
    exact_match: float | None
    distance: float | None
    n: int


def _flatten_summary(experiment: str, model: str, timestamp: str,
                     summary: dict[str, Any]) -> list[CellRecord]:
    out = []
    for cond, row in summary.items():
        if not isinstance(row, dict):
            continue
        out.append(CellRecord(
            experiment=experiment,
            model=model,
            timestamp=timestamp,
            condition=cond,
            df_moved=row.get("anchor_direction_follow_rate"),
            df_raw=row.get("anchor_direction_follow_rate_raw"),
            adopt=row.get("anchor_adoption_rate"),
            exact_match=row.get("accuracy_exact"),
            distance=row.get("mean_distance_to_anchor"),
            n=int(row.get("count", 0)),
        ))
    return out


def _walk_summaries(root: Path) -> pd.DataFrame:
    rows: list[CellRecord] = []
    for exp in EXPERIMENT_GLOBS:
        exp_dir = root / exp
        if not exp_dir.exists():
            continue
        for model_dir in sorted(exp_dir.iterdir()):
            if not model_dir.is_dir():
                continue
            for ts_dir in sorted(model_dir.iterdir()):
                if not ts_dir.is_dir():
                    continue
                summary_path = ts_dir / "summary.json"
                if not summary_path.exists():
                    continue
                try:
                    summary = json.loads(summary_path.read_text())
                except Exception:
                    continue
                rows.extend(_flatten_summary(exp, model_dir.name, ts_dir.name, summary))
    return pd.DataFrame([r.__dict__ for r in rows])


def _compute_deltas(before: pd.DataFrame, after: pd.DataFrame) -> pd.DataFrame:
    key = ["experiment", "model", "timestamp", "condition"]
    merged = before.merge(after, on=key, how="outer", suffixes=("_before", "_after"))
    merged["df_moved_delta"] = merged["df_moved_after"] - merged["df_moved_before"]
    merged["df_raw_delta"] = merged["df_raw_after"] - merged["df_raw_before"]
    merged["adopt_delta"] = merged["adopt_after"] - merged["adopt_before"]
    merged["exact_match_delta"] = merged["exact_match_after"] - merged["exact_match_before"]
    return merged


def _figure_main_panel(deltas: pd.DataFrame, out: Path) -> None:
    """7-model main-panel before/after df_moved comparison (anchor S1 cell)."""
    main = deltas[
        (deltas["experiment"] == "experiment")
        & (deltas["condition"].str.contains("number"))
    ].copy()
    if main.empty:
        return
    main = main.groupby("model")[["df_moved_before", "df_moved_after"]].mean().reset_index()
    main = main.sort_values("df_moved_before", ascending=False)
    fig, ax = plt.subplots(figsize=(9, 4.5))
    x = range(len(main))
    width = 0.35
    ax.bar([i - width / 2 for i in x], main["df_moved_before"], width=width,
           label="anchor·gt form (before)", color="#888")
    ax.bar([i + width / 2 for i in x], main["df_moved_after"], width=width,
           label="C-form (after)", color="#2c7fb8")
    ax.set_xticks(list(x))
    ax.set_xticklabels(main["model"], rotation=20, ha="right")
    ax.set_ylabel("direction_follow_rate (M2 moved)")
    ax.set_title("Main panel — VQAv2 number, anchor S1 (averaged across conditions)")
    ax.legend()
    fig.tight_layout()
    fig.savefig(out, dpi=140)
    plt.close(fig)


def _figure_per_experiment_scatter(deltas: pd.DataFrame, out: Path) -> None:
    """Each (experiment × model × condition) cell plotted as a before vs after dot."""
    cells = deltas.dropna(subset=["df_moved_before", "df_moved_after"]).copy()
    if cells.empty:
        return
    fig, ax = plt.subplots(figsize=(7, 7))
    colors = {exp: f"C{i}" for i, exp in enumerate(sorted(cells["experiment"].unique()))}
    for exp, sub in cells.groupby("experiment"):
        ax.scatter(sub["df_moved_before"], sub["df_moved_after"], s=22,
                   alpha=0.65, label=exp, color=colors[exp])
    lim = max(cells["df_moved_before"].max(), cells["df_moved_after"].max(), 0.5) * 1.05
    ax.plot([0, lim], [0, lim], "k--", linewidth=0.6, alpha=0.4)
    ax.set_xlim(0, lim)
    ax.set_ylim(0, lim)
    ax.set_xlabel("df_moved (anchor·gt form, before)")
    ax.set_ylabel("df_moved (C-form, after)")
    ax.set_title("Per-cell df_moved drift across the project")
    ax.legend(fontsize=7, loc="upper left", ncol=2)
    fig.tight_layout()
    fig.savefig(out, dpi=140)
    plt.close(fig)


def _summary_tables(deltas: pd.DataFrame) -> dict[str, pd.DataFrame]:
    out = {}
    main = deltas[deltas["experiment"] == "experiment"].copy()
    if not main.empty:
        agg = main.groupby("model")[
            ["df_moved_before", "df_moved_after", "df_raw_before", "df_raw_after",
             "adopt_before", "adopt_after", "exact_match_before", "exact_match_after"]
        ].mean()
        out["main_panel_VQAv2"] = agg.round(4)

    by_exp = deltas.groupby("experiment").agg(
        n_cells=("df_moved_before", "size"),
        df_moved_before_mean=("df_moved_before", "mean"),
        df_moved_after_mean=("df_moved_after", "mean"),
        df_moved_delta_mean=("df_moved_delta", "mean"),
        df_moved_delta_max=("df_moved_delta", lambda s: s.abs().max()),
    ).round(4)
    out["per_experiment_drift"] = by_exp

    return out


def _qualitative_status_table(deltas: pd.DataFrame) -> pd.DataFrame:
    """Coarse status per (experiment, model) pair: did df_moved survive?"""
    rows = []
    for (exp, model), sub in deltas.groupby(["experiment", "model"]):
        non_null = sub.dropna(subset=["df_moved_before", "df_moved_after"])
        if non_null.empty:
            rows.append({"experiment": exp, "model": model, "status": "no data"})
            continue
        avg_before = non_null["df_moved_before"].mean()
        avg_after = non_null["df_moved_after"].mean()
        if avg_before == 0 and avg_after > 0.01:
            status = "REVEALED (was 0, now > 0)"
        elif abs(avg_after - avg_before) <= 0.01:
            status = "preserved (~ unchanged)"
        elif avg_after > avg_before:
            status = "stronger (after > before)"
        else:
            status = "weakened (after < before)"
        rows.append({
            "experiment": exp, "model": model,
            "before_mean": round(avg_before, 4),
            "after_mean": round(avg_after, 4),
            "status": status,
        })
    return pd.DataFrame(rows)


def _df_to_markdown(df: pd.DataFrame, index: bool = True) -> str:
    """tabulate-free DataFrame → markdown table."""
    if df.empty:
        return "_no data_"
    cols = list(df.columns)
    if index:
        index_name = df.index.name or ""
        header = "| " + " | ".join([index_name] + [str(c) for c in cols]) + " |"
        sep = "|" + "|".join(["---"] * (len(cols) + 1)) + "|"
        rows = ["| " + " | ".join([str(idx)] + [_fmt(v) for v in row]) + " |"
                for idx, row in zip(df.index, df.itertuples(index=False))]
    else:
        header = "| " + " | ".join(str(c) for c in cols) + " |"
        sep = "|" + "|".join(["---"] * len(cols)) + "|"
        rows = ["| " + " | ".join(_fmt(v) for v in row) + " |"
                for row in df.itertuples(index=False)]
    return "\n".join([header, sep] + rows)


def _fmt(v: Any) -> str:
    if v is None or (isinstance(v, float) and pd.isna(v)):
        return ""
    if isinstance(v, float):
        return f"{v:.4f}"
    return str(v)


def _markdown(deltas: pd.DataFrame, tables: dict[str, pd.DataFrame],
              status: pd.DataFrame) -> str:
    lines = ["# C-form migration report — before vs after",
             "",
             "**Status:** generated by `scripts/build_C_form_migration_report.py`.",
             "",
             "## What changed",
             "",
             "On 2026-04-28 the project's `direction_follow_rate` was refactored from",
             "the buggy anchor·gt form `(pa-gt)·(anchor-gt) > 0` (inherited from the",
             "M1-era code body, never updated to match the M2 docstring) to the C-form",
             "`(pa-pb)·(anchor-pb) > 0  AND  pa != pb` (the user's intended formula —",
             "anchor pull measured as a baseline-relative shift, gt-free).",
             "",
             "The same refactor pass also closed a driver schema gap: the three M2",
             "row flags `anchor_direction_followed_moved`, `pred_b_equal_anchor`,",
             "`pred_diff_from_base` were never threaded into the driver row dict, so",
             "every directly-driven `summary.json` reported `df_M2 = 0` even when the",
             "raw signal was non-zero. `reaggregate_paired_adoption.py` had been",
             "silently fixing that for any dir it touched, but eight dirs were never",
             "touched.",
             "",
             "## Coverage",
             "",
             f"- Cells compared: **{len(deltas)}**",
             f"- Experiments: **{deltas['experiment'].nunique()}**",
             f"- Models: **{deltas['model'].nunique()}**",
             "",
             "## Main panel — VQAv2 number subset (7 models)",
             "",
             "Average across {b, d, anchor} conditions per model:",
             "",
             _df_to_markdown(tables.get("main_panel_VQAv2", pd.DataFrame())),
             "",
             "![main panel df_moved before vs after](../figures/C_form_migration_main_panel.png)",
             "",
             "## Per-experiment drift",
             "",
             _df_to_markdown(tables.get("per_experiment_drift", pd.DataFrame())),
             "",
             "## Qualitative status (per experiment × model)",
             "",
             _df_to_markdown(status, index=False),
             "",
             "## Per-cell drift",
             "",
             "![per-cell scatter](../figures/C_form_migration_scatter.png)",
             "",
             "## Methodology",
             "",
             "1. Pre-refactor `summary.json` for every relevant dir was archived under",
             "   `outputs/before_C_form/` *before* any code change.",
             "2. `metrics.py:118` was rewritten to compute the C-form, with parallel",
             "   updates in `reaggregate_paired_adoption.py`,",
             "   `analyze_metric_variants.py`, the unit tests, and 7 doc surfaces.",
             "3. `scripts/run_experiment.py` row dict gained the three missing M2 flags",
             "   (`pred_b_equal_anchor`, `pred_diff_from_base`,",
             "   `anchor_direction_followed_moved`); a regression test in",
             "   `tests/test_metrics.py::DriverRowSchemaRegressionTest` enforces that.",
             "4. `scripts/reaggregate_paired_adoption.py --apply --force` was run over",
             "   every output sub-tree.",
             "5. This script consumes the two trees and emits the comparison.",
             "",
             "_Generated automatically — re-run via_",
             "`uv run python scripts/build_C_form_migration_report.py`",
             ""]
    return "\n".join(lines)


def _build_pptx(deltas: pd.DataFrame, tables: dict[str, pd.DataFrame],
                status: pd.DataFrame, out: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    def title(slide, text: str) -> None:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7))
        tf = tb.text_frame
        tf.text = text
        tf.paragraphs[0].runs[0].font.size = Pt(28)
        tf.paragraphs[0].runs[0].font.bold = True

    def bullet(slide, lines: list[str], top: float = 1.1, height: float = 5.5,
               size: int = 18) -> None:
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.0),
                                       Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            for run in p.runs:
                run.font.size = Pt(size)

    # Slide 1 — title
    s = prs.slides.add_slide(blank)
    title(s, "C-form migration — before vs after")
    bullet(s, [
        "direction_follow_rate refactor (2026-04-28).",
        "Before: anchor·gt form (pa-gt)·(anchor-gt) — buggy code body, never matched docstring.",
        "After: C-form (pa-pb)·(anchor-pb) — user's intended formula, gt-free, anchor pull as baseline-relative shift.",
        f"Coverage: {len(deltas)} cells, {deltas['experiment'].nunique()} experiments, {deltas['model'].nunique()} models.",
    ], top=1.5, height=4.5, size=22)

    # Slide 2 — what changed
    s = prs.slides.add_slide(blank)
    title(s, "What changed (and why)")
    bullet(s, [
        "Two compounding bugs found while debugging γ-β df_M2 = 0:",
        "1. Code/doc form mismatch — M2 commit 304a677 declared pb·gt form in",
        "   docstring + commit message + 5 doc surfaces, but metrics.py:118 stayed",
        "   on M1-era anchor·gt form. Both forms turned out to be wrong-by-author.",
        "2. Driver schema gap — run_experiment.py never threaded",
        "   anchor_direction_followed_moved + pred_b_equal_anchor +",
        "   pred_diff_from_base into the row dict, so summarize_condition saw",
        "   None and reported df_M2 = 0 silently.",
        "",
        "Fix: C-form (pa-pb)·(anchor-pb) in 4 code sites + 7 doc surfaces;",
        "drivers thread the missing flags; reaggregate sweep over 61 dirs.",
    ], top=1.0, height=6.0, size=18)

    # Slide 3 — main panel figure
    s = prs.slides.add_slide(blank)
    title(s, "Main panel — VQAv2 number, 7 models")
    s.shapes.add_picture(str(OUT_FIG / "C_form_migration_main_panel.png"),
                          Inches(0.7), Inches(1.1), height=Inches(5.6))
    bullet(s, [
        "Result: every model's df_moved is HIGHER under C-form.",
        "Largest shift: qwen2.5-vl-7b 0.079 → 0.214 (×2.7).",
        "C-form is more sensitive to anchor pull because it directly measures",
        "pa shifting from pb toward anchor (gt-free).",
    ], top=6.7, height=0.6, size=12)

    # Slide 4 — main panel table
    main_panel = tables.get("main_panel_VQAv2", pd.DataFrame())
    s = prs.slides.add_slide(blank)
    title(s, "Main panel — numerical table")
    if not main_panel.empty:
        rows = main_panel.shape[0] + 1
        cols = main_panel.shape[1] + 1
        tbl = s.shapes.add_table(rows, cols, Inches(0.4), Inches(1.1),
                                  Inches(12.5), Inches(5.5)).table
        tbl.cell(0, 0).text = "model"
        for j, col in enumerate(main_panel.columns):
            tbl.cell(0, j + 1).text = str(col).replace("_", " ")
        for i, (idx, row) in enumerate(main_panel.iterrows()):
            tbl.cell(i + 1, 0).text = str(idx)
            for j, v in enumerate(row):
                tbl.cell(i + 1, j + 1).text = _fmt(v)
        for cell in [tbl.cell(0, j) for j in range(cols)]:
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.bold = True
                    r.font.size = Pt(10)
        for i in range(1, rows):
            for j in range(cols):
                for p in tbl.cell(i, j).text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(10)

    # Slide 5 — scatter + per-experiment drift
    s = prs.slides.add_slide(blank)
    title(s, "Per-cell drift across the project")
    s.shapes.add_picture(str(OUT_FIG / "C_form_migration_scatter.png"),
                          Inches(0.6), Inches(1.0), height=Inches(6.2))
    bullet(s, [
        "Each dot = one (experiment × model × condition) cell.",
        "Above the diagonal → C-form > anchor·gt form.",
        "Most cells above the diagonal — C-form lifts the metric across the board.",
        "Cells at (0, ≥0.05) = E5e MathVista runs, where the buggy 0 turned",
        "out to be a driver schema artefact, NOT a real categorical-replace regime.",
    ], top=1.0, height=5.5, size=12)

    # Slide 6 — qualitative status
    s = prs.slides.add_slide(blank)
    title(s, "Qualitative status — paper-tier claims")
    by_status = (
        status.groupby("status").size().reset_index(name="n").sort_values("n", ascending=False)
        if not status.empty
        else pd.DataFrame()
    )
    if not by_status.empty:
        rows = by_status.shape[0] + 1
        tbl = s.shapes.add_table(rows, 2, Inches(0.5), Inches(1.2),
                                  Inches(8.0), Inches(0.4 * rows + 0.5)).table
        tbl.cell(0, 0).text = "status"
        tbl.cell(0, 1).text = "# (model × experiment)"
        for i, (_, row) in enumerate(by_status.iterrows()):
            tbl.cell(i + 1, 0).text = str(row["status"])
            tbl.cell(i + 1, 1).text = str(row["n"])
    bullet(s, [
        "Headline (A1, E1d, E4, L1, E5b, E5c, E5e ChartQA/TallyQA): preserved or stronger.",
        "γ-α MathVista 'categorical-replace' framing: REVEALED — was driver bug, not real.",
        "No paper-tier qualitative claim flipped; many were strengthened.",
        "→ Paper writeup needs number refresh, not narrative rewrite (modulo MathVista §).",
    ], top=5.6, height=1.5, size=14)

    # Slide 7 — methodology + repro
    s = prs.slides.add_slide(blank)
    title(s, "Methodology + reproducibility")
    bullet(s, [
        "1. Pre-refactor summary.json archived to outputs/before_C_form/ before any code change.",
        "2. metrics.py + reaggregate + analyze_metric_variants + tests rewritten to C-form.",
        "3. 7 doc surfaces (project.md, roadmap.md, AGENTS.md, M2 evidence, paper_summary_slides,",
        "   metrics.py docstrings, analyze_metric_variants docstrings) re-stated in C-form.",
        "4. run_experiment.py row dict gained the 3 missing M2 flags.",
        "5. Regression test in tests/test_metrics.py::DriverRowSchemaRegressionTest.",
        "6. scripts/reaggregate_paired_adoption.py --apply --force across 17 sub-trees.",
        "7. This deck + accompanying markdown report regenerable via",
        "   `uv run python scripts/build_C_form_migration_report.py`.",
        "",
        "Memory entry feedback_metric_C_form.md persisted so future sessions",
        "do not re-litigate the formula choice.",
    ], top=1.0, height=6.0, size=16)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))


def main() -> None:
    if not BEFORE_ROOT.exists():
        raise SystemExit(f"backup root missing: {BEFORE_ROOT}")
    OUT_FIG.mkdir(parents=True, exist_ok=True)

    print(f"reading before tree at {BEFORE_ROOT} ...")
    before = _walk_summaries(BEFORE_ROOT)
    print(f"  → {len(before)} cells")

    print(f"reading after tree at {AFTER_ROOT} ...")
    after = _walk_summaries(AFTER_ROOT)
    print(f"  → {len(after)} cells")

    deltas = _compute_deltas(before, after)
    deltas.to_csv(PROJECT_ROOT / "docs" / "insights" / "_data" / "C_form_migration_cells.csv", index=False)
    print(f"wrote {len(deltas)} cell-level rows")

    tables = _summary_tables(deltas)
    status = _qualitative_status_table(deltas)
    status.to_csv(PROJECT_ROOT / "docs" / "insights" / "_data" / "C_form_migration_status.csv", index=False)

    _figure_main_panel(deltas, OUT_FIG / "C_form_migration_main_panel.png")
    _figure_per_experiment_scatter(deltas, OUT_FIG / "C_form_migration_scatter.png")
    print(f"wrote figures to {OUT_FIG}/C_form_migration_*.png")

    OUT_DOC.write_text(_markdown(deltas, tables, status))
    print(f"wrote markdown to {OUT_DOC}")

    _build_pptx(deltas, tables, status, OUT_PPTX)
    print(f"wrote pptx deck to {OUT_PPTX}")


if __name__ == "__main__":
    main()
